"""文本提取与切块

从多种格式文档中提取纯文本，并按 token 切分为适合向量化的块。
支持格式：TXT、MD、PDF、DOCX、XLSX、CSV、PPTX、PPT

智能 chunk 策略流程：
1. 分析文档特征（长度、密度、复杂度等）
2. 根据分析结果制定个性化 chunk 策略
3. 执行切块操作
"""

import csv
import logging
import os
import re
from typing import List, Optional

from models import Chunk

logger = logging.getLogger(__name__)


class TextChunker:
    """文本提取与切块器

    负责从多种格式的文档中提取纯文本，并将文本切分为适合向量化的块。
    
    智能切块策略流程：
    1. 分析文档特征（长度、密度、复杂度等）
    2. 根据分析结果制定个性化 chunk 策略
    3. 执行切块操作，按中文句子边界优先切分
    """

    SUPPORTED_EXTENSIONS = {
        ".txt": "_extract_txt",
        ".md": "_extract_markdown",
        ".docx": "_extract_word",
        ".xlsx": "_extract_excel",
        ".pdf": "_extract_pdf",
        ".csv": "_extract_csv",
        ".pptx": "_extract_ppt",
        ".ppt": "_extract_ppt",
    }

    _CHUNK_SIZE_BY_TYPE = {
        ".txt": 512,
        ".md": 512,
        ".docx": 512,
        ".xlsx": 256,
        ".pdf": 384,
        ".csv": 256,
        ".pptx": 512,
        ".ppt": 512,
    }

    _OVERLAP_BY_TYPE = {
        ".txt": 64,
        ".md": 64,
        ".docx": 64,
        ".xlsx": 32,
        ".pdf": 48,
        ".csv": 32,
        ".pptx": 64,
        ".ppt": 64,
    }

    _SENTENCE_BOUNDARIES = set("。！？；\n")

    def __init__(self, tokenizer=None, chunk_size: int = None, chunk_overlap: int = None):
        """初始化 TextChunker

        Args:
            tokenizer: tokenizer 实例（用于 token 计数和切块）
            chunk_size: 目标块大小（token），为 None 时按策略自动选择
            chunk_overlap: 重叠区域大小（token），为 None 时按策略自动选择
        """
        self.tokenizer = tokenizer
        self.default_chunk_size = chunk_size
        self.default_chunk_overlap = chunk_overlap

    def _get_chunk_params(self, file_path: str) -> tuple:
        """根据文件类型获取基础 chunk 参数

        Args:
            file_path: 文件路径

        Returns:
            (chunk_size, chunk_overlap) 元组
        """
        if self.default_chunk_size is not None and self.default_chunk_overlap is not None:
            return (self.default_chunk_size, self.default_chunk_overlap)

        ext = os.path.splitext(file_path)[1].lower()
        chunk_size = self._CHUNK_SIZE_BY_TYPE.get(ext, 512)
        chunk_overlap = self._OVERLAP_BY_TYPE.get(ext, 64)
        return (chunk_size, chunk_overlap)

    def analyze_document(self, text: str, file_path: str = None) -> dict:
        """分析文档特征，为制定 chunk 策略提供依据

        Args:
            text: 原始文本
            file_path: 文件路径

        Returns:
            文档分析结果字典，包含：
            - total_chars: 总字符数
            - total_tokens: 总 token 数
            - avg_sentence_length: 平均句子长度（字符）
            - paragraph_count: 段落数量
            - density_score: 文本密度分数（0-1）
            - complexity_level: 复杂度级别（low/medium/high）
            - title_density: 标题密度（标题行数/总行数）
            - list_density: 列表密度（列表行数/总行数）
            - has_tables: 是否包含表格结构
            - content_type: 内容类型（narrative/list/table/mixed）
            - avg_paragraph_length: 平均段落长度
        """
        if not text or not text.strip():
            return {
                "total_chars": 0,
                "total_tokens": 0,
                "avg_sentence_length": 0,
                "paragraph_count": 0,
                "density_score": 0,
                "complexity_level": "low",
                "title_density": 0,
                "list_density": 0,
                "has_tables": False,
                "content_type": "narrative",
                "avg_paragraph_length": 0,
            }

        total_chars = len(text)

        if self.tokenizer:
            total_tokens = len(self.tokenizer.encode(text, add_special_tokens=False))
        else:
            total_tokens = int(total_chars / 2.5)

        lines = text.split("\n")
        total_lines = len([l for l in lines if l.strip()])
        
        paragraphs = [p for p in text.split("\n\n") if p.strip()]
        paragraph_count = len(paragraphs)

        sentences = re.split(r"[。！？；]", text)
        sentences = [s.strip() for s in sentences if s.strip()]
        avg_sentence_length = sum(len(s) for s in sentences) / max(len(sentences), 1)

        non_whitespace_chars = sum(1 for c in text if not c.isspace())
        density_score = non_whitespace_chars / max(total_chars, 1)

        title_lines = 0
        list_lines = 0
        has_tables = False
        
        for line in lines:
            stripped = line.strip()
            if stripped:
                if stripped.startswith(("#", "##", "###", "####", "#####", "######", "【", "】", "第", "一、", "二、", "三、", "1.", "2.", "3.", "（一）", "（二）")):
                    title_lines += 1
                elif stripped.startswith(("* ", "- ", "• ", "· ", "○ ", "● ", "□ ", "■ ", "1.", "2.", "3.", "（1）", "（2）", "（3）")):
                    list_lines += 1
                elif "\t" in line or "|" in line and stripped.count("|") >= 2:
                    has_tables = True

        title_density = title_lines / max(total_lines, 1)
        list_density = list_lines / max(total_lines, 1)

        avg_paragraph_length = sum(len(p) for p in paragraphs) / max(paragraph_count, 1)

        if avg_sentence_length > 100 or density_score > 0.9:
            complexity_level = "high"
        elif avg_sentence_length > 50 or density_score > 0.7:
            complexity_level = "medium"
        else:
            complexity_level = "low"

        if list_density > 0.4:
            content_type = "list"
        elif has_tables or "\t" in text and total_chars > 1000:
            content_type = "table"
        elif title_density > 0.2:
            content_type = "structured"
        else:
            content_type = "narrative"

        logger.info(f"[文档分析] 文件: {os.path.basename(file_path) if file_path else '未知'}")
        logger.info(f"  总字符数: {total_chars:,}, 总Token数: {total_tokens:,}")
        logger.info(f"  段落数: {paragraph_count}, 平均句子长度: {avg_sentence_length:.1f}")
        logger.info(f"  复杂度: {complexity_level}, 密度分数: {density_score:.2f}")
        logger.info(f"  标题密度: {title_density:.2f}, 列表密度: {list_density:.2f}")
        logger.info(f"  是否含表格: {has_tables}, 内容类型: {content_type}")

        return {
            "total_chars": total_chars,
            "total_tokens": total_tokens,
            "avg_sentence_length": avg_sentence_length,
            "paragraph_count": paragraph_count,
            "density_score": density_score,
            "complexity_level": complexity_level,
            "title_density": title_density,
            "list_density": list_density,
            "has_tables": has_tables,
            "content_type": content_type,
            "avg_paragraph_length": avg_paragraph_length,
        }

    def determine_chunk_strategy(self, analysis: dict, file_path: str = None) -> dict:
        """根据文档分析结果制定个性化的 chunk 策略

        Args:
            analysis: 文档分析结果
            file_path: 文件路径

        Returns:
            chunk 策略字典，包含：
            - chunk_size: 块大小（token）
            - chunk_overlap: 重叠区域大小（token）
            - min_chunk_size: 最小块大小（token）
            - merge_short_chunks: 是否合并小块
            - boundary_preference: 边界偏好（sentence/paragraph/none）
            - strategy_name: 策略名称（用于日志和调试）
        """
        base_chunk_size, base_overlap = self._get_chunk_params(file_path or "")

        complexity = analysis.get("complexity_level", "medium")
        avg_sentence_length = analysis.get("avg_sentence_length", 50)
        content_type = analysis.get("content_type", "narrative")
        title_density = analysis.get("title_density", 0)
        list_density = analysis.get("list_density", 0)
        has_tables = analysis.get("has_tables", False)

        chunk_size = base_chunk_size
        chunk_overlap = base_overlap
        strategy_name = "默认策略"

        if content_type == "list":
            chunk_size = int(base_chunk_size * 0.8)
            chunk_overlap = int(base_overlap * 1.2)
            strategy_name = "列表型文档策略"
        elif content_type == "table":
            chunk_size = int(base_chunk_size * 0.5)
            chunk_overlap = int(base_overlap * 1.5)
            strategy_name = "表格型文档策略"
        elif content_type == "structured":
            chunk_size = int(base_chunk_size * 0.9)
            chunk_overlap = int(base_overlap * 1.1)
            strategy_name = "结构化文档策略"
        elif complexity == "high":
            chunk_size = int(base_chunk_size * 0.75)
            chunk_overlap = int(base_overlap * 1.5)
            strategy_name = "高复杂度策略"
        elif complexity == "low":
            chunk_size = int(base_chunk_size * 1.25)
            chunk_overlap = int(base_overlap * 0.75)
            strategy_name = "低复杂度策略"

        if avg_sentence_length > 120:
            chunk_size = min(int(chunk_size * 1.3), 1024)
            strategy_name += "+长句适配"
        elif avg_sentence_length < 30:
            chunk_size = int(chunk_size * 0.85)
            strategy_name += "+短句适配"

        if title_density > 0.15:
            chunk_overlap = int(chunk_overlap * 1.2)
            strategy_name += "+标题增强"

        if list_density > 0.2:
            chunk_size = int(chunk_size * 0.9)
            chunk_overlap = int(chunk_overlap * 1.1)

        min_chunk_size = max(64, int(chunk_size * 0.3))

        if analysis.get("paragraph_count", 0) > 15:
            boundary_preference = "paragraph"
        elif has_tables or content_type == "table":
            boundary_preference = "paragraph"
        elif avg_sentence_length < 40:
            boundary_preference = "sentence"
        else:
            boundary_preference = "sentence"

        strategy = {
            "chunk_size": int(chunk_size),
            "chunk_overlap": int(chunk_overlap),
            "min_chunk_size": int(min_chunk_size),
            "merge_short_chunks": content_type != "list",
            "boundary_preference": boundary_preference,
            "strategy_name": strategy_name,
        }

        logger.info(f"[Chunk策略] 策略名称: {strategy_name}")
        logger.info(f"  Chunk大小: {chunk_size} token, 重叠: {chunk_overlap} token")
        logger.info(f"  最小块大小: {min_chunk_size} token, 边界偏好: {boundary_preference}")
        logger.info(f"  合并小块: {strategy['merge_short_chunks']}")

        return strategy

    def chunk(self, text: str, file_path: str = None) -> List[Chunk]:
        """将文本切分为多个 Chunk

        执行流程：
        1. 分析文档特征（字符数、token数、句子长度、段落数、密度、复杂度等）
        2. 根据分析结果制定个性化 chunk 策略（块大小、重叠、边界偏好）
        3. 执行智能切块（优先按中文句子/段落边界切分）

        Args:
            text: 原始文本
            file_path: 文件路径（用于确定 chunk size）

        Returns:
            切块列表
        """
        if not self.tokenizer:
            raise ValueError("tokenizer 未设置，无法进行切块")

        if not text or not text.strip():
            return []

        file_name = os.path.basename(file_path) if file_path else "未知文件"
        logger.info(f"\n[Chunk处理开始] 文件: {file_name}")

        analysis = self.analyze_document(text, file_path)
        logger.info(f"[分析完成] 准备制定策略")

        strategy = self.determine_chunk_strategy(analysis, file_path)
        chunk_size = strategy["chunk_size"]
        chunk_overlap = strategy["chunk_overlap"]
        min_chunk_size = strategy["min_chunk_size"]
        boundary_preference = strategy["boundary_preference"]
        logger.info(f"[策略确定] 开始执行切块")

        token_ids = self.tokenizer.encode(text, add_special_tokens=False)
        total_tokens = len(token_ids)

        if total_tokens == 0:
            logger.warning(f"[Chunk处理结束] 文本为空，返回空列表")
            return []

        if total_tokens <= chunk_size:
            logger.info(f"[Chunk处理结束] 文本较短，无需切块（{total_tokens} token）")
            return [Chunk(text=text.strip(), index=0, token_count=total_tokens)]

        text_len = len(text)
        avg_token_len = text_len / max(total_tokens, 1)

        chunks: List[Chunk] = []
        start = 0
        chunk_index = 0

        while start < total_tokens:
            end = min(int(start + chunk_size), total_tokens)

            if end < total_tokens:
                if boundary_preference == "paragraph":
                    boundary = self._find_paragraph_boundary(text, int(start), int(end), avg_token_len)
                else:
                    boundary = self._find_sentence_boundary(token_ids, int(start), int(end), int(chunk_overlap))
                if boundary > start and boundary <= end:
                    end = boundary
                    logger.debug(f"  边界调整: {end - chunk_overlap} -> {end}")

            chunk_token_ids = token_ids[int(start):int(end)]
            chunk_token_count = len(chunk_token_ids)

            char_start = int(start * avg_token_len)
            char_end = min(int(end * avg_token_len), text_len)

            if char_end > char_start and char_end <= text_len:
                chunk_text = text[char_start:char_end]
            else:
                chunk_text = self.tokenizer.decode(chunk_token_ids, skip_special_tokens=True)

            chunk_text = chunk_text.strip()
            if chunk_text:
                chunks.append(Chunk(
                    text=chunk_text,
                    index=chunk_index,
                    token_count=chunk_token_count,
                ))
                logger.debug(f"  Chunk {chunk_index}: {chunk_token_count} token, {len(chunk_text)} 字符")

            if end >= total_tokens:
                break

            start = end - chunk_overlap
            chunk_index += 1

        if strategy["merge_short_chunks"] and len(chunks) > 1:
            merged_count = len(chunks)
            chunks = self._merge_short_chunks(chunks, min_chunk_size)
            logger.info(f"[合并小块] {merged_count} -> {len(chunks)} 个")

        logger.info(f"[Chunk处理完成] 共生成 {len(chunks)} 个文本块")

        return chunks

    def _find_paragraph_boundary(self, text: str, start_token: int, end_token: int, avg_token_len: float) -> int:
        """在文本中找到段落边界"""
        char_start = int(start_token * avg_token_len)
        char_end = min(int(end_token * avg_token_len), len(text))
        
        search_start = max(char_start, int((start_token + end_token * 0.5) * avg_token_len))
        
        range_start = min(char_end, len(text) - 1)
        range_stop = search_start - 1
        
        if range_start <= range_stop:
            return end_token
        
        for i in range(range_start, range_stop, -1):
            if text[i] == '\n' and (i + 1 >= len(text) or text[i + 1] == '\n'):
                return int(i / avg_token_len) + 1
            if text[i] == '\n':
                return int(i / avg_token_len) + 1
        
        return end_token

    def _merge_short_chunks(self, chunks: List[Chunk], min_chunk_size: int) -> List[Chunk]:
        """合并小于最小尺寸的相邻块"""
        if len(chunks) <= 1:
            return chunks

        merged = []
        current = chunks[0]

        for next_chunk in chunks[1:]:
            if current.token_count < min_chunk_size:
                current = Chunk(
                    text=current.text + "\n" + next_chunk.text,
                    index=current.index,
                    token_count=current.token_count + next_chunk.token_count,
                )
            else:
                merged.append(current)
                current = next_chunk

        merged.append(current)
        return merged

    def _get_chunk_start(self, chunks: List[Chunk], chunk_overlap: int) -> int:
        """计算最后一个 chunk 的起始 token 位置"""
        if len(chunks) <= 1:
            return 0
        start = 0
        for chunk in chunks[:-1]:
            end = start + chunk.token_count
            start = end - chunk_overlap
        return start

    def _find_sentence_boundary(self, token_ids: list, start: int, max_end: int, chunk_overlap: int) -> int:
        """在 token 序列中找到最近的中文句子边界

        从 max_end 向前搜索，找到最近的句子结束标点位置。
        """
        search_start = max(start + chunk_overlap, start)

        for pos in range(max_end - 1, search_start - 1, -1):
            token_text = self.tokenizer.decode([token_ids[pos]], skip_special_tokens=True)
            for ch in token_text:
                if ch in self._SENTENCE_BOUNDARIES:
                    return pos + 1

        return max_end

    def extract_text(self, file_path: str) -> Optional[str]:
        """从文件提取纯文本

        根据文件扩展名分发到对应的提取器。
        提取失败或内容为空时返回 None。
        所有提取的文本都会经过中文空格清理。

        Args:
            file_path: 文件绝对路径

        Returns:
            提取的纯文本内容，失败返回 None
        """
        ext = os.path.splitext(file_path)[1].lower()

        if ext not in self.SUPPORTED_EXTENSIONS:
            logger.warning("不支持的文件格式: %s", file_path)
            return None

        method_name = self.SUPPORTED_EXTENSIONS[ext]
        extract_method = getattr(self, method_name)

        try:
            text = extract_method(file_path)
        except Exception as e:
            logger.error("文件提取失败 [%s]: %s", file_path, str(e))
            return None

        if text is None or text.strip() == "":
            logger.warning("文件内容为空: %s", file_path)
            return None

        text = self._normalize_chinese_spacing(text)
        
        return text

    def _extract_txt(self, file_path: str) -> Optional[str]:
        """提取 TXT 文件（UTF-8）"""
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    def _extract_markdown(self, file_path: str) -> Optional[str]:
        """提取 Markdown 文件，移除语法标记保留纯文本"""
        from markdown_it import MarkdownIt

        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()

        md = MarkdownIt()
        tokens = md.parse(content)

        text_parts = []
        for token in tokens:
            if token.type == "inline" and token.children:
                inline_text = self._extract_inline_text(token.children)
                if inline_text:
                    text_parts.append(inline_text)
            elif token.type == "fence":
                if token.content:
                    text_parts.append(token.content.rstrip("\n"))

        return "\n".join(text_parts)

    def _extract_inline_text(self, children) -> str:
        """从 inline token 的 children 中提取纯文本"""
        parts = []
        for child in children:
            if child.type == "text":
                parts.append(child.content)
            elif child.type == "code_inline":
                parts.append(child.content)
            elif child.type in ("softbreak", "hardbreak"):
                parts.append("\n")
        return "".join(parts)

    def _extract_word(self, file_path: str) -> Optional[str]:
        """提取 Word (.docx) 文件，按段落提取"""
        from docx import Document

        try:
            doc = Document(file_path)
            paragraphs = []
            for para in doc.paragraphs:
                try:
                    text = para.text
                    if text and text.strip():
                        paragraphs.append(text.strip())
                except Exception:
                    continue
            return "\n".join(paragraphs) if paragraphs else None
        except Exception as e:
            logger.error(f"Word文档解析失败 [{file_path}]: {e}")
            return None

    def _extract_excel(self, file_path: str) -> Optional[str]:
        """提取 Excel (.xlsx) 文件，按行制表符分隔"""
        from openpyxl import load_workbook

        wb = load_workbook(file_path, read_only=True, data_only=True)
        all_text = []

        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                cells = []
                for cell in row:
                    value = cell.value
                    if value is not None:
                        cells.append(str(value))
                    else:
                        cells.append("")
                all_text.append("\t".join(cells))

        wb.close()
        return "\n".join(all_text)

    def _extract_pdf(self, file_path: str) -> Optional[str]:
        """提取 PDF 文件，按页顺序拼接"""
        from PyPDF2 import PdfReader

        try:
            reader = PdfReader(file_path)
            pages_text = []

            for page in reader.pages:
                try:
                    text = page.extract_text()
                    if text:
                        text = self._normalize_chinese_spacing(text)
                        if text.strip():
                            pages_text.append(text)
                except Exception as e:
                    logger.warning(f"PDF页面提取失败 [{file_path}]: {e}")
                    continue

            return "\n".join(pages_text) if pages_text else None
        except Exception as e:
            logger.error(f"PDF解析失败 [{file_path}]: {e}")
            return None

    def _normalize_chinese_spacing(self, text: str) -> str:
        """清理中文文本中的多余空格

        处理的空白字符：普通空格(0x20)、非断行空格(0xA0)、全角空格(0x3000)
        """
        chinese_char = r"[\u4e00-\u9fff]"
        chinese_punctuation = r"[\u3002\uff0c\uff01\uff1f\uff1b\uff1a\u201c\u201d\u2018\u2019\u300a\u300b\u3008\u3009\u3010\u3011\u3001\u00b7]"

        space_chars = r" \xa0\u3000"

        pattern1 = re.compile(f"({chinese_char})[{space_chars}]+({chinese_char})")
        pattern2 = re.compile(f"({chinese_char})[{space_chars}]+({chinese_punctuation})")
        pattern3 = re.compile(f"({chinese_punctuation})[{space_chars}]+({chinese_char})")

        prev_text = text
        while True:
            text = pattern1.sub(r"\1\2", text)
            text = pattern2.sub(r"\1\2", text)
            text = pattern3.sub(r"\1\2", text)
            if text == prev_text:
                break
            prev_text = text

        return text

    def _extract_csv(self, file_path: str) -> Optional[str]:
        """提取 CSV 文件（UTF-8），按行制表符分隔"""
        rows = []
        with open(file_path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append("\t".join(row))
        return "\n".join(rows)

    def _extract_ppt(self, file_path: str) -> Optional[str]:
        """提取 PowerPoint (.pptx) 文件，按幻灯片顺序拼接文本"""
        from pptx import Presentation
        from pptx.exc import PackageNotFoundError

        try:
            prs = Presentation(file_path)
            slides_text = []

            for i, slide in enumerate(prs.slides, 1):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                if slide_content:
                    slides_text.append(f"[幻灯片 {i}]")
                    slides_text.extend(slide_content)

            return "\n\n".join(slides_text)

        except PackageNotFoundError:
            logger.warning(f"文件格式可能不是有效的 .pptx，仅支持现代 PowerPoint 格式: {file_path}")
            return None
        except Exception as e:
            logger.error(f"PowerPoint 提取失败 [{file_path}]: {e}")
            return None